from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Slide 1 : Titre
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Architecture Microservices"
slide.placeholders[1].text = "Vue d'ensemble mise a jour"

# Slide 2 : Schéma visuel
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Schema d'Architecture"
img_path = "architecture.png"
slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))

# Slide 3 : Légende des composants
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Composants et Connexions"
tf = slide.shapes.placeholders[1].text_frame
tf.clear()
entries = [
    ("API Gateway", "Entrée unique, routage, JWT, CORS"),
    ("Auth Service", "Generation des Token JWT"),
    ("Eureka Server", "Découverte dynamique des microservices"),
    #("Spring Cloud Config", "Config centralisée, secrets Vault"),
    #("Prometheus/Grafana/Loki", "Monitoring métriques and Logs")
]
for title, desc in entries:
    p = tf.add_paragraph()
    p.text = f"{title} : {desc}"
    p.level = 0

# Slide 4 : Bonus sécurités & flows JWT
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Team 1 Service"
tf = slide.shapes.placeholders[1].text_frame
tf.clear()
for line in [
    "Auth Service",
    "MongoDB",
    "Exécute tests via Selenium"
]:
    p = tf.add_paragraph()
    p.text = f"✔ {line}"
    p.level = 0

# Slide 5 : Bonus sécurités & flows JWT
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Team 2 Service"
tf = slide.shapes.placeholders[1].text_frame
tf.clear()
for line in [
    "Auth Service",
    "MongoDB",
    "Exécute tests via Selenium"
]:
    p = tf.add_paragraph()
    p.text = f"✔ {line}"
    p.level = 0


# Slide 5 : Bonus sécurités & flows JWT
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Team 2 Service"
tf = slide.shapes.placeholders[1].text_frame
tf.clear()
for line in [
    "Auth Service",
    "MongoDB",
    "Mysql",
    "Exécute tests via Selenium"
]:
    p = tf.add_paragraph()
    p.text = f"✔ {line}"
    p.level = 0

# Slide 6 : Bonus sécurités & flows JWT
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Sécurité & Flows JWT"
tf = slide.shapes.placeholders[1].text_frame
tf.clear()
for line in [
    "Rotation de clés JWT via JWKS",
    "Health checks & métriques Actuator"
]:
    p = tf.add_paragraph()
    p.text = f"✔ {line}"
    p.level = 0

prs.save("architecture_presentation.pptx")
print("Présentation générée : architecture_presentation.pptx")
